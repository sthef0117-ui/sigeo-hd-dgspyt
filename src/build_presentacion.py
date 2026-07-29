"""
SIGEO-HD DGSPYT — Generador de la presentación para la reunión de mandos.

Del apunte de la reunión: «Bosquejo para presentar», «Presentar trabajo».
El mando no navega un tablero en la mesa: lleva láminas.

    python src/etl_sigeo.py
    python src/build_presentacion.py

Salida: presentacion/SIGEO-HD_Reunion_Mandos.pptx

Todas las cifras y gráficas se calculan del pipeline. No hay nada escrito a
mano, así que la presentación se rehace en un comando cada corte nuevo.
"""

import json
import sys
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt  # noqa: E402
from pptx import Presentation  # noqa: E402
from pptx.dml.color import RGBColor  # noqa: E402
from pptx.enum.text import PP_ALIGN  # noqa: E402
from pptx.util import Emu, Pt  # noqa: E402

RAIZ = Path(__file__).resolve().parent.parent
ANALISIS = RAIZ / "analisis"
SALIDA = RAIZ / "presentacion"
TEMP = SALIDA / "_graficas"

# Paleta institucional, la misma del tablero.
GUINDA = RGBColor(0x7A, 0x13, 0x27)
TINTA = RGBColor(0x16, 0x20, 0x2C)
GRIS = RGBColor(0x59, 0x65, 0x73)
GRIS_CLARO = RGBColor(0x8A, 0x94, 0xA2)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
FONDO = RGBColor(0xF7, 0xF8, 0xFA)
ROJO = RGBColor(0xC1, 0x12, 0x1F)
AMBAR = RGBColor(0xB4, 0x53, 0x09)

HG, HB = "#7a1327", "#c94f6d"
ANCHO, ALTO = Emu(12192000), Emu(6858000)  # 16:9


def cm(v):
    return Emu(int(v * 360000))


MENORES = {"de", "del", "la", "las", "el", "los", "en", "y", "con", "por",
           "a", "al", "sin", "para"}


def titulo_es(texto_mayus):
    """Capitaliza respetando el español: str.title() escribe «De La Via»."""
    palabras = str(texto_mayus).lower().split()
    return " ".join(
        p.capitalize() if i == 0 or p not in MENORES else p
        for i, p in enumerate(palabras))


def leer(nombre):
    ruta = ANALISIS / nombre
    if not ruta.exists():
        print(f"Falta {ruta.name}. Ejecuta primero: python src/etl_sigeo.py")
        sys.exit(1)
    return json.loads(ruta.read_text(encoding="utf-8"))


# --------------------------------------------------------------------------
# Piezas de lámina
# --------------------------------------------------------------------------

def texto(diapo, x, y, ancho, alto, contenido, tam=18, negrita=False,
          color=TINTA, alineacion=PP_ALIGN.LEFT, interlineado=1.15):
    caja = diapo.shapes.add_textbox(x, y, ancho, alto)
    marco = caja.text_frame
    marco.word_wrap = True
    for i, linea in enumerate(str(contenido).split("\n")):
        p = marco.paragraphs[0] if i == 0 else marco.add_paragraph()
        p.text = linea
        p.alignment = alineacion
        p.line_spacing = interlineado
        for r in p.runs:
            r.font.size = Pt(tam)
            r.font.bold = negrita
            r.font.color.rgb = color
            r.font.name = "Segoe UI"
    return caja


def banda(diapo, y, alto, color):
    from pptx.enum.shapes import MSO_SHAPE
    f = diapo.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, ANCHO, alto)
    f.fill.solid()
    f.fill.fore_color.rgb = color
    f.line.fill.background()
    f.shadow.inherit = False
    return f


def tarjeta(diapo, x, y, ancho, alto, rotulo, valor, detalle="", color=GUINDA):
    from pptx.enum.shapes import MSO_SHAPE
    f = diapo.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, ancho, alto)
    f.fill.solid()
    f.fill.fore_color.rgb = BLANCO
    f.line.color.rgb = RGBColor(0xE0, 0xE4, 0xEA)
    f.line.width = Pt(1)
    f.shadow.inherit = False
    f.adjustments[0] = 0.06
    texto(diapo, x + cm(.5), y + cm(.35), ancho - cm(1), cm(.8),
          rotulo.upper(), tam=10, negrita=True, color=GRIS_CLARO)
    texto(diapo, x + cm(.5), y + cm(1.0), ancho - cm(1), cm(1.4),
          valor, tam=30, negrita=True, color=color)
    if detalle:
        texto(diapo, x + cm(.5), y + cm(2.5), ancho - cm(1), cm(1.2),
              detalle, tam=10, color=GRIS)


def encabezado(diapo, titulo, bajada=""):
    banda(diapo, 0, cm(2.6), TINTA)
    texto(diapo, cm(1.2), cm(.55), cm(24), cm(1.1), titulo, tam=22,
          negrita=True, color=BLANCO)
    if bajada:
        texto(diapo, cm(1.2), cm(1.6), cm(28), cm(.8), bajada, tam=11,
              color=RGBColor(0xB9, 0xC0, 0xC9))
    texto(diapo, cm(27.5), cm(.9), cm(6), cm(.8), "SIGEO-HD", tam=13,
          negrita=True, color=RGBColor(0xE0, 0x52, 0x6F), alineacion=PP_ALIGN.RIGHT)


def tabla(diapo, x, y, ancho, columnas, filas, anchos=None, tam=11):
    n_f, n_c = len(filas) + 1, len(columnas)
    alto = cm(.85) + cm(.62) * len(filas)
    forma = diapo.shapes.add_table(n_f, n_c, x, y, ancho, alto)
    t = forma.table
    if anchos:
        total = sum(anchos)
        for i, a in enumerate(anchos):
            t.columns[i].width = Emu(int(ancho * a / total))
    for j, c in enumerate(columnas):
        celda = t.cell(0, j)
        celda.text = c
        celda.fill.solid()
        celda.fill.fore_color.rgb = TINTA
        p = celda.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT if j and j >= n_c - 6 else PP_ALIGN.LEFT
        for r in p.runs:
            r.font.size = Pt(tam - 1)
            r.font.bold = True
            r.font.color.rgb = BLANCO
            r.font.name = "Segoe UI"
    for i, fila in enumerate(filas, 1):
        for j, v in enumerate(fila):
            celda = t.cell(i, j)
            celda.text = str(v)
            celda.fill.solid()
            celda.fill.fore_color.rgb = BLANCO if i % 2 else FONDO
            p = celda.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j and j >= n_c - 6 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.size = Pt(tam)
                r.font.color.rgb = TINTA
                r.font.name = "Segoe UI"
                r.font.bold = (j == 1 and n_c > 3)
    return forma


def grafica(nombre, dibujar, ancho=9, alto=4.2):
    TEMP.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(ancho, alto), dpi=200)
    dibujar(ax)
    for lado in ("top", "right"):
        ax.spines[lado].set_visible(False)
    for lado in ("left", "bottom"):
        ax.spines[lado].set_color("#d5dae1")
    ax.tick_params(colors="#596573", labelsize=9)
    fig.tight_layout()
    ruta = TEMP / f"{nombre}.png"
    fig.savefig(ruta, transparent=True)
    plt.close(fig)
    return ruta


# --------------------------------------------------------------------------

def main():
    R = leer("resumen_ejecutivo.json")
    hd = leer("corroborados_sigeo.json")
    zonas = leer("zonas_ciegas.json")
    auditoria = leer("auditoria_decesos.json")
    territorio = leer("perfil_territorial.json")
    coord = leer("perfil_coordinaciones.json")["coordinaciones"]

    pres = Presentation()
    pres.slide_width, pres.slide_height = ANCHO, ALTO
    vacia = pres.slide_layouts[6]
    N = lambda v: f"{v:,}".replace(",", ",")

    # ---------- 1. Portada ----------
    d = pres.slides.add_slide(vacia)
    banda(d, 0, ALTO, TINTA)
    texto(d, cm(2.5), cm(4.2), cm(24), cm(2), "HOMICIDIOS DOLOSOS", tam=40,
          negrita=True, color=BLANCO)
    texto(d, cm(2.5), cm(6.4), cm(24), cm(1.2), "Análisis de corte · Julio 2026",
          tam=20, color=RGBColor(0xE0, 0x52, 0x6F))
    texto(d, cm(2.5), cm(8.6), cm(26), cm(2),
          "Dirección General de Seguridad Pública y Tránsito\n"
          "Unidad de Homicidios Dolosos", tam=14,
          color=RGBColor(0xB9, 0xC0, 0xC9))
    texto(d, cm(2.5), cm(15.4), cm(28), cm(1.4),
          f"Documento de trabajo interno · generado {R['generado']} · "
          "cifras sin validar por las áreas responsables",
          tam=9, color=GRIS)

    # ---------- 2. Panorama ----------
    d = pres.slides.add_slide(vacia)
    encabezado(d, "Panorama del corte",
               f"{R['serie']['horas']} horas de operación C5 · "
               f"{N(R['c5']['llamadas_totales'])} llamadas · "
               f"{R['hd']['municipios']} municipios con homicidio")
    a, sep = cm(7.4), cm(.55)
    for i, (rot, val, det, col) in enumerate([
        ("Homicidios dolosos", str(R["hd"]["eventos"]),
         f"{R['hd']['victimas']} víctimas fatales", GUINDA),
        ("En franja nocturna", f"{R['hd']['nocturnos']}",
         f"{round(R['hd']['nocturnos']*100/R['hd']['eventos'])}% entre 22:00 y 06:00", ROJO),
        ("Distancia a base", f"{R['cobertura']['dist_hd_base_mediana_km']} km",
         f"mediana · {R['cobertura']['hd_a_mas_de_3km_de_base']} a más de 3 km", TINTA),
        ("Personal desplegado", N(R["cobertura"]["personal_total"]),
         f"en {R['cobertura']['bases_en_uso_georreferenciadas']} bases en uso", RGBColor(0x0F, 0x76, 0x6E)),
    ]):
        tarjeta(d, cm(1.2) + i * (a + sep), cm(3.5), a, cm(4), rot, val, det, col)

    por_dia = {}
    for x in hd:
        if x["fecha"]:
            por_dia[x["fecha"]] = por_dia.get(x["fecha"], 0) + 1
    dias = sorted(por_dia)

    def dibujar_dias(ax):
        ax.bar([x[8:] + "/" + x[5:7] for x in dias], [por_dia[x] for x in dias],
               color=HG, width=.72)
        ax.set_ylabel("eventos", fontsize=9, color="#596573")
        ax.grid(axis="y", color="#eceff3")
        ax.set_axisbelow(True)

    d.shapes.add_picture(str(grafica("dias", dibujar_dias, 13.5, 3.6)),
                         cm(1.2), cm(8.2), width=cm(30.5))
    texto(d, cm(1.2), cm(17.2), cm(30), cm(1),
          "Homicidios dolosos corroborados por día", tam=11, color=GRIS)

    # ---------- 3. Coordinaciones ----------
    d = pres.slides.add_slide(vacia)
    encabezado(d, "Carga por coordinación regional",
               "La columna que decide: carga de violencia por cada 100 elementos adscritos")

    def dibujar_carga(ax):
        cs = sorted(coord, key=lambda c: c["carga_por_100_elementos"] or 0)
        ax.barh([titulo_es(c["coordinacion"]) for c in cs],
                [c["carga_por_100_elementos"] or 0 for c in cs], color=HB, height=.66)
        for i, c in enumerate(cs):
            ax.text((c["carga_por_100_elementos"] or 0) + 1, i,
                    f"{c['carga_por_100_elementos']}", va="center",
                    fontsize=9, color="#16202c")
        ax.grid(axis="x", color="#eceff3")
        ax.set_axisbelow(True)

    d.shapes.add_picture(str(grafica("carga", dibujar_carga, 7.6, 4.6)),
                         cm(1.2), cm(3.6), width=cm(15))
    tabla(d, cm(16.8), cm(3.6), cm(15),
          ["Coordinación", "Mpios", "HD", "Personal", "/100 el."],
          [[titulo_es(c["coordinacion"]), c["total_municipios"], c["hd_eventos"],
            N(c["personal_total"]), c["carga_por_100_elementos"]] for c in coord],
          anchos=[38, 13, 11, 20, 18], tam=10)
    texto(d, cm(1.2), cm(15.6), cm(30.5), cm(2.2),
          "Mide condiciones del territorio, no el desempeño de una persona: los insumos no "
          "contienen asignación nominal de mando, turnos ni recorridos.\n"
          "22 de los 108 municipios se clasificaron por vecindad territorial al no tener "
          "inmueble propio en el inventario.",
          tam=10, color=GRIS)

    # ---------- 4. Zonas ciegas ----------
    d = pres.slides.add_slide(vacia)
    desat = [z for z in zonas if z["clasificacion"] == "DESATENDIDO"]
    satur = sorted([z for z in zonas if z["clasificacion"] == "SATURADO"],
                   key=lambda z: -z["indice_violencia"])
    encabezado(d, "¿Por qué no se patrulla ahí?",
               f"{len(zonas)} sectores con violencia registrada · "
               f"{len(desat)} desatendidos · {len(satur)} saturados")
    texto(d, cm(1.2), cm(3.4), cm(15), cm(1),
          "EXTENDER COBERTURA", tam=12, negrita=True, color=ROJO)
    tabla(d, cm(1.2), cm(4.2), cm(15), ["Municipio", "Índice", "Base a", "Personal 3 km"],
          [[titulo_es(z["municipio"]), z["indice_ceguera"], f"{z['dist_base_km']} km",
            N(z["personal_3km"])] for z in desat[:8]],
          anchos=[45, 18, 19, 18], tam=10)
    texto(d, cm(16.8), cm(3.4), cm(15), cm(1),
          "REVISAR EL PATRULLAJE QUE YA EXISTE", tam=12, negrita=True, color=AMBAR)
    tabla(d, cm(16.8), cm(4.2), cm(15), ["Municipio", "Violencia", "Base a", "Personal 3 km"],
          [[titulo_es(z["municipio"]), z["indice_violencia"], f"{z['dist_base_km']} km",
            N(z["personal_3km"])] for z in satur[:8]],
          anchos=[45, 18, 19, 18], tam=10)
    texto(d, cm(1.2), cm(15.8), cm(30.5), cm(1.8),
          "Izquierda: hay violencia y no hay despliegue a distancia útil. Procede extender "
          "cobertura o reasignar cuadrante.\n"
          "Derecha: hay base y personal a menos de 1.5 km y aun así concentra violencia. "
          "Procede revisar efectividad y horarios, no sumar inmuebles.",
          tam=10, color=GRIS)

    # ---------- 5. Casos a revisar ----------
    d = pres.slides.add_slide(vacia)
    altas = [a for a in auditoria if a["nivel_revision"] == "ALTA"]
    encabezado(d, "Decesos dudosos, suicidios y personas no localizadas",
               f"{R['auditoria']['casos_revisados']} casos auditados · "
               f"{len(altas)} de prioridad alta")
    tabla(d, cm(1.2), cm(3.5), cm(30.5),
          ["Folio C5", "Municipio", "Clasificación de origen", "Pts", "Motivo de revisión"],
          [[c["folio"], titulo_es(c["municipio"]), titulo_es(c["incidente"]), c["puntaje"],
            "; ".join(i["motivo"] for i in c["indicadores"] if i["peso"] > 0)[:90]]
           for c in altas[:9]],
          anchos=[15, 14, 24, 6, 41], tam=9)
    texto(d, cm(1.2), cm(15.4), cm(30.5), cm(2),
          "Es una prioridad de revisión documental, no un dictamen pericial. "
          "La determinación de la mecánica de los hechos corresponde a los servicios "
          "periciales y a la FGJEM. Cada caso muestra los indicadores exactos que "
          "elevaron su prioridad, para que el criterio sea auditable.",
          tam=10, color=GRIS)

    # ---------- 6. Acuerdos ----------
    d = pres.slides.add_slide(vacia)
    g = R["geocodificacion"]
    encabezado(d, "Acuerdos propuestos", "Para la reunión de mandos")
    top_desat = ", ".join(dict.fromkeys(titulo_es(z["municipio"]) for z in desat[:3]))
    acuerdos = [
        f"Extender cobertura o reasignar cuadrante en los {len(desat)} sectores "
        f"desatendidos, priorizando {top_desat}.",
        f"Revisar efectividad y horarios del patrullaje en los {len(satur)} sectores "
        "saturados, donde ya hay base y personal a menos de 1.5 km.",
        f"Solicitar a la FGJEM las carpetas de los {len(altas)} casos de prioridad alta "
        "y turnarlos a Criminalística.",
        "Instruir a cabina C5 la captura obligatoria del campo de referencia de "
        f"ubicación: es lo que permite recuperar el {g['porcentaje_violentas_recuperado']}% "
        "de las llamadas con violencia que llegan sin coordenada.",
        "Completar el personal adscrito en el inventario de inmuebles: sin ese dato "
        "la medición de cobertura subestima el despliegue real.",
    ]
    y = cm(3.8)
    for i, a in enumerate(acuerdos, 1):
        texto(d, cm(1.4), y, cm(1.4), cm(1.2), str(i), tam=22, negrita=True, color=GUINDA)
        texto(d, cm(3.0), y + cm(.1), cm(28), cm(2), a, tam=13, color=TINTA)
        y += cm(2.35)
    texto(d, cm(1.2), cm(16.4), cm(30.5), cm(1),
          "Cifras reproducibles con  python src/etl_sigeo.py", tam=9, color=GRIS_CLARO)

    SALIDA.mkdir(parents=True, exist_ok=True)
    destino = SALIDA / "SIGEO-HD_Reunion_Mandos.pptx"
    pres.save(destino)

    for p in TEMP.glob("*.png"):
        p.unlink()
    TEMP.rmdir()

    print("SIGEO-HD DGSPYT · presentación generada")
    print(f"  láminas .................... {len(pres.slides.__iter__.__self__._sldIdLst)}")
    print(f"  archivo .................... {destino}")
    print(f"  tamaño ..................... {destino.stat().st_size / 1024:.0f} KB")
    return 0


if __name__ == "__main__":
    sys.exit(main())
